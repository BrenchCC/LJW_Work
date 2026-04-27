import os
import logging
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(15, 23, 42)
TEXT = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
BLUE = RGBColor(30, 64, 175)
LIGHT = RGBColor(244, 247, 255)
BORDER = RGBColor(226, 232, 240)
GOLD = RGBColor(180, 126, 34)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(22, 101, 52)
RED = RGBColor(185, 28, 28)


def text_box(
    slide,
    text,
    left,
    top,
    width,
    height,
    size,
    color = TEXT,
    bold = False,
    align = PP_ALIGN.LEFT
):
    """Add a single-paragraph text box.

    Parameters:
        slide: Slide object to update.
        text: Text content.
        left: Horizontal position.
        top: Vertical position.
        width: Text box width.
        height: Text box height.
        size: Font size in points.
        color: Text color.
        bold: Whether text should be bold.
        align: Paragraph alignment.
    """
    shape = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = "PingFang SC"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return shape


def bullets(slide, items, left, top, width, height, size = 16):
    """Add a readable bullet list.

    Parameters:
        slide: Slide object to update.
        items: Bullet strings.
        left: Horizontal position.
        top: Vertical position.
        width: Text box width.
        height: Text box height.
        size: Font size in points.
    """
    shape = slide.shapes.add_textbox(
        left,
        top,
        width,
        height
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        paragraph.font.name = "PingFang SC"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = TEXT
    return shape


def title(slide, main, sub = None):
    """Add a consistent title area.

    Parameters:
        slide: Slide object to update.
        main: Main title.
        sub: Optional subtitle.
    """
    text_box(
        slide,
        main,
        Inches(0.62),
        Inches(0.32),
        Inches(12.1),
        Inches(0.58),
        27,
        NAVY,
        True
    )
    line = slide.shapes.add_shape(
        1,
        Inches(0.64),
        Inches(1.02),
        Inches(12.05),
        Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.color.rgb = BLUE
    if sub:
        text_box(
            slide,
            sub,
            Inches(0.66),
            Inches(1.11),
            Inches(11.8),
            Inches(0.3),
            11,
            MUTED
        )


def card(slide, left, top, width, height, heading, body, accent = BLUE, body_size = 13):
    """Add a light card with an accent bar.

    Parameters:
        slide: Slide object to update.
        left: Horizontal position.
        top: Vertical position.
        width: Card width.
        height: Card height.
        heading: Card heading.
        body: Card body.
        accent: Accent color.
        body_size: Body font size.
    """
    rect = slide.shapes.add_shape(
        5,
        left,
        top,
        width,
        height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT
    rect.line.color.rgb = BORDER
    bar = slide.shapes.add_shape(
        1,
        left,
        top,
        Inches(0.06),
        height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    text_box(
        slide,
        heading,
        left + Inches(0.18),
        top + Inches(0.12),
        width - Inches(0.35),
        Inches(0.28),
        14,
        accent,
        True
    )
    text_box(
        slide,
        body,
        left + Inches(0.18),
        top + Inches(0.48),
        width - Inches(0.34),
        height - Inches(0.58),
        body_size,
        TEXT
    )


def image_contain(slide, path, left, top, width, height):
    """Add an image without distorting its aspect ratio.

    Parameters:
        slide: Slide object to update.
        path: Image path.
        left: Target box horizontal position.
        top: Target box vertical position.
        width: Target box width.
        height: Target box height.
    """
    if not path.exists():
        logger.warning("Missing image: %s", path)
        return None
    with Image.open(path) as img:
        img_w, img_h = img.size
    scale = min(width / img_w, height / img_h)
    draw_w = int(img_w * scale)
    draw_h = int(img_h * scale)
    draw_left = left + int((width - draw_w) / 2)
    draw_top = top + int((height - draw_h) / 2)
    return slide.shapes.add_picture(
        str(path),
        draw_left,
        draw_top,
        width = draw_w,
        height = draw_h
    )


def footer(slide, number):
    """Add footer metadata.

    Parameters:
        slide: Slide object to update.
        number: Slide number.
    """
    text_box(
        slide,
        f"SubspaceAD / Feature Subspace Projection  ·  {number:02d}",
        Inches(0.66),
        Inches(7.08),
        Inches(6.0),
        Inches(0.22),
        8,
        MUTED
    )


def blank(prs):
    """Create one blank slide.

    Parameters:
        prs: Presentation object.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def flow(slide, labels, left, top, width, color = BLUE):
    """Draw a horizontal process flow.

    Parameters:
        slide: Slide object to update.
        labels: Box labels.
        left: Horizontal start.
        top: Vertical start.
        width: Total flow width.
        color: Accent color.
    """
    gap = Inches(0.18)
    box_w = int((width - gap * (len(labels) - 1)) / len(labels))
    for index, label in enumerate(labels):
        x = left + index * (box_w + gap)
        card(
            slide,
            x,
            top,
            box_w,
            Inches(0.9),
            f"Step {index + 1}",
            label,
            color,
            12
        )
        if index < len(labels) - 1:
            text_box(
                slide,
                "->",
                x + box_w + Inches(0.02),
                top + Inches(0.31),
                Inches(0.18),
                Inches(0.18),
                12,
                MUTED,
                True,
                PP_ALIGN.CENTER
            )


def build(output_path):
    """Build the enriched presentation.

    Parameters:
        output_path: Destination PPTX path.
    """
    root = Path(os.getcwd())
    img = root / "2026-04-27" / "results" / "assay_analysis" / "images"
    bg = root / "2026-04-27" / "results" / "background_pages"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = blank(prs)
    text_box(
        slide,
        "SubspaceAD",
        Inches(0.75),
        Inches(0.95),
        Inches(6.1),
        Inches(0.72),
        43,
        NAVY,
        True
    )
    text_box(
        slide,
        "方法解析与业务落地",
        Inches(0.78),
        Inches(1.76),
        Inches(6.0),
        Inches(0.55),
        30,
        BLUE,
        True
    )
    text_box(
        slide,
        "Training-Free Few-Shot Anomaly Detection via Subspace Modeling",
        Inches(0.82),
        Inches(2.55),
        Inches(5.8),
        Inches(0.34),
        14,
        TEXT
    )
    card(
        slide,
        Inches(0.82),
        Inches(4.82),
        Inches(5.7),
        Inches(1.18),
        "任务定位",
        "从论文中提炼可验证 baseline 与实验问题，判断它能否服务新电池零样本缺陷检测中的特征子空间投影任务。",
        BLUE,
        15
    )
    image_contain(
        slide,
        img / "header.png",
        Inches(6.65),
        Inches(0.95),
        Inches(5.9),
        Inches(1.75)
    )
    image_contain(
        slide,
        img / "technical_framework.png",
        Inches(6.55),
        Inches(3.08),
        Inches(5.95),
        Inches(2.35)
    )
    footer(slide, 1)

    slide = blank(prs)
    title(slide, "业务背景：为什么不能直接复用旧缺陷模型", "新型号上线时缺陷样本没有标签，domain shift 会破坏直接复用")
    image_contain(
        slide,
        bg / "page-03.png",
        Inches(0.72),
        Inches(1.52),
        Inches(5.85),
        Inches(3.25)
    )
    card(slide, Inches(6.95), Inches(1.55), Inches(5.45), Inches(1.0), "可用资源", "旧型号：有划痕样本与标签；新型号：只有产线采集图像。")
    card(slide, Inches(6.95), Inches(2.83), Inches(5.45), Inches(1.0), "困难", "缺陷物理本质相似，但材质、光照、反光率、背景和纹理改变了成像域。", GOLD)
    card(slide, Inches(6.95), Inches(4.11), Inches(5.45), Inches(1.0), "目标", "先验证预训练特征中是否存在可投影、可迁移的缺陷/正常子空间结构。", GREEN)
    bullets(
        slide,
        [
            "传统微调：重新采集 -> 标注 -> 训练，周期长。",
            "子空间路线：先把判别问题搬到特征空间，再看投影后的距离/残差。",
        ],
        Inches(0.88),
        Inches(5.15),
        Inches(11.5),
        Inches(1.0),
        15
    )
    footer(slide, 2)

    slide = blank(prs)
    title(slide, "论文定位：少样本正常图像能否替代复杂训练", "不训练网络、不做 prompt tuning、不维护 memory bank")
    card(slide, Inches(0.8), Inches(1.65), Inches(3.55), Inches(1.18), "输入设定", "每类只给 k = 1, 2, 4 张正常训练图像；测试图像可能正常或异常。")
    card(slide, Inches(4.88), Inches(1.65), Inches(3.55), Inches(1.18), "核心假设", "DINOv2 patch-level 表征足够强，异常判断可由简单统计模型完成。", GREEN)
    card(slide, Inches(8.96), Inches(1.65), Inches(3.55), Inches(1.18), "输出", "图像级异常分数 + 像素级异常定位热图。")
    text_box(slide, "一句话方法", Inches(0.92), Inches(3.42), Inches(2.0), Inches(0.3), 16, BLUE, True)
    text_box(
        slide,
        "用正常 patch 特征拟合 PCA 子空间；测试 patch 到该子空间的重构残差就是异常分数。",
        Inches(0.92),
        Inches(3.82),
        Inches(11.2),
        Inches(0.65),
        24,
        NAVY,
        True
    )
    flow(
        slide,
        [
            "Frozen DINOv2\npatch features",
            "PCA normal\nsubspace",
            "Projection\nresidual",
            "Anomaly map\n+ image score",
        ],
        Inches(0.92),
        Inches(5.25),
        Inches(11.3)
    )
    footer(slide, 3)

    slide = blank(prs)
    title(slide, "方法总览：两阶段闭环", "训练阶段只拟合 PCA；推理阶段只做投影和残差计算")
    image_contain(
        slide,
        img / "technical_framework.png",
        Inches(0.82),
        Inches(1.48),
        Inches(11.72),
        Inches(3.0)
    )
    card(slide, Inches(0.95), Inches(4.78), Inches(3.55), Inches(1.12), "Fitting", "正常图像 -> DINOv2 patch features -> PCA 子空间参数 mu, C。")
    card(slide, Inches(4.9), Inches(4.78), Inches(3.55), Inches(1.12), "Inference", "测试图像 -> 同一 backbone -> 投影到正常子空间。")
    card(slide, Inches(8.85), Inches(4.78), Inches(3.55), Inches(1.12), "Scoring", "残差图用于定位；top 1% patch 分数用于图像级检测。")
    footer(slide, 4)

    slide = blank(prs)
    title(slide, "特征抽取细节：为什么用中间层和增强", "这部分决定 PCA 协方差估计是否稳定")
    card(slide, Inches(0.8), Inches(1.55), Inches(3.75), Inches(1.15), "Backbone / Resolution", "DINOv2-G, ViT-G/14；输入统一为 672 x 672。")
    card(slide, Inches(0.8), Inches(2.95), Inches(3.75), Inches(1.15), "Layer aggregation", "聚合第 22-28 层中间特征，兼顾局部结构与语义信息。", GREEN)
    card(slide, Inches(0.8), Inches(4.35), Inches(3.75), Inches(1.15), "Augmentation", "每张正常图生成 Na = 30 个旋转增强视图，覆盖工业场景中的姿态变化。", GOLD)
    image_contain(
        slide,
        img / "figure_4_resolution.png",
        Inches(5.0),
        Inches(1.47),
        Inches(3.45),
        Inches(3.55)
    )
    image_contain(
        slide,
        img / "figure_5_backbone.png",
        Inches(8.85),
        Inches(1.47),
        Inches(3.45),
        Inches(3.55)
    )
    text_box(slide, "工程含义：子空间方法本身很轻，主要成本由 backbone 和输入分辨率决定。", Inches(5.05), Inches(5.42), Inches(7.0), Inches(0.45), 15, TEXT, True)
    footer(slide, 5)

    slide = blank(prs)
    title(slide, "PCA 正常子空间：论文的核心建模", "正常变化被看作低维线性子空间，异常位于残差方向")
    image_contain(
        slide,
        img / "technical_pca_scoring.png",
        Inches(0.75),
        Inches(1.42),
        Inches(3.0),
        Inches(5.05)
    )
    card(slide, Inches(4.25), Inches(1.48), Inches(3.75), Inches(1.15), "Normal features", "X_normal 收集 k 张正常图及增强视图的所有 patch 特征。")
    card(slide, Inches(8.35), Inches(1.48), Inches(3.75), Inches(1.15), "PCA parameters", "计算均值 mu 和协方差 Sigma；C 为前 r 个特征向量。")
    text_box(slide, "x = mu + C z + epsilon", Inches(4.45), Inches(3.25), Inches(7.1), Inches(0.52), 28, BLUE, True, PP_ALIGN.CENTER)
    text_box(slide, "sum(lambda_i, i <= r) >= tau * sum(lambda_i),  tau = 0.99", Inches(4.45), Inches(4.05), Inches(7.1), Inches(0.42), 17, NAVY, True, PP_ALIGN.CENTER)
    bullets(
        slide,
        [
            "C 表示正常外观主变化方向。",
            "tau 不能取 1.00，因为必须保留残差空间来承载异常信号。",
        ],
        Inches(4.45),
        Inches(4.82),
        Inches(7.2),
        Inches(1.2),
        15
    )
    footer(slide, 6)

    slide = blank(prs)
    title(slide, "异常评分：投影残差、图像级聚合、像素级定位", "从 patch residual 到 heatmap，再到 image score")
    flow(
        slide,
        [
            "Extract test\npatch feature x_p",
            "Project:\nx_proj = mu + CC^T(x_p-mu)",
            "Residual:\nS = ||x_p-x_proj||^2",
            "Upsample +\nGaussian smooth",
            "TVaR top 1%\nimage score",
        ],
        Inches(0.75),
        Inches(1.55),
        Inches(11.9)
    )
    card(slide, Inches(0.82), Inches(3.2), Inches(5.65), Inches(1.35), "为什么 top 1% 聚合", "稀疏划痕可能只占少量 patch；平均全图会被大量正常背景稀释，top tail 更敏感。", GREEN)
    card(slide, Inches(6.9), Inches(3.2), Inches(5.65), Inches(1.35), "为什么保留热图", "业务上不仅要二分类，还要给质检人员提供可复核证据：缺陷区域、阈值依据、误检来源。", GOLD)
    image_contain(
        slide,
        img / "figure_3_qualitative.png",
        Inches(1.0),
        Inches(4.92),
        Inches(11.35),
        Inches(1.55)
    )
    footer(slide, 7)

    slide = blank(prs)
    title(slide, "主实验：少样本设置下达到强基线", "MVTec-AD 与 VisA，1/2/4-shot 正常样本")
    image_contain(
        slide,
        img / "table_1_main_comparison.png",
        Inches(0.7),
        Inches(1.38),
        Inches(7.2),
        Inches(4.9)
    )
    card(slide, Inches(8.32), Inches(1.55), Inches(3.95), Inches(0.95), "MVTec-AD 1-shot", "I-AUROC 97.1%\nP-AUROC 97.5%", GREEN, 16)
    card(slide, Inches(8.32), Inches(2.76), Inches(3.95), Inches(0.95), "VisA 1-shot", "I-AUROC 93.4%\nP-AUROC 98.2%", GREEN, 16)
    card(slide, Inches(8.32), Inches(3.97), Inches(3.95), Inches(1.25), "相较 AnomalyDINO", "VisA 1-shot I-AUROC +6.0；PRO +1.0。", BLUE, 15)
    card(slide, Inches(8.32), Inches(5.48), Inches(3.95), Inches(0.82), "谨慎结论", "多数指标领先，但不是所有单项都绝对第一。", GOLD, 13)
    footer(slide, 8)

    slide = blank(prs)
    title(slide, "定性结果：残差热图可以定位缺陷", "热图质量决定它能否用于人工复核和阈值调参")
    image_contain(
        slide,
        img / "figure_3_qualitative.png",
        Inches(0.82),
        Inches(1.45),
        Inches(11.8),
        Inches(4.8)
    )
    text_box(slide, "业务启发：即使最终交付是“划痕 / 正常”，也应保存 top patch、热图和阈值证据，方便质检人员复核。", Inches(1.0), Inches(6.35), Inches(11.2), Inches(0.42), 15, NAVY, True)
    footer(slide, 9)

    slide = blank(prs)
    title(slide, "消融结论：哪些因素真正关键", "子空间要足够表达正常变化，但不能吞掉异常残差")
    image_contain(slide, img / "figure_4_resolution.png", Inches(0.75), Inches(1.55), Inches(3.2), Inches(2.65))
    image_contain(slide, img / "figure_5_backbone.png", Inches(4.15), Inches(1.55), Inches(3.2), Inches(2.65))
    image_contain(slide, img / "table_4_pca_variance.png", Inches(7.55), Inches(1.55), Inches(3.2), Inches(2.65))
    card(slide, Inches(0.85), Inches(4.72), Inches(3.15), Inches(1.25), "Resolution", "448 以上整体趋稳；672 是跨数据集稳健折中。")
    card(slide, Inches(4.25), Inches(4.72), Inches(3.15), Inches(1.25), "Backbone", "DINOv2-G 最强；边缘部署可用 S/B/L 做速度折中。")
    card(slide, Inches(7.65), Inches(4.72), Inches(3.15), Inches(1.25), "PCA tau", "0.95-0.99 稳；1.00 明显下降，说明残差空间很关键。", GOLD)
    card(slide, Inches(10.95), Inches(1.62), Inches(1.55), Inches(4.35), "Takeaway", "强特征决定上限；PCA 阈值决定异常信号是否被保留下来。", GREEN, 12)
    footer(slide, 10)

    slide = blank(prs)
    title(slide, "论文方法如何迁移到业务任务", "SubspaceAD 是 Normal-PCA baseline，不等于最终业务方案")
    card(slide, Inches(0.82), Inches(1.55), Inches(3.7), Inches(1.25), "可直接借鉴", "特征抽取、PCA 子空间、投影残差、TVaR 聚合、热图可视化。", GREEN)
    card(slide, Inches(4.85), Inches(1.55), Inches(3.7), Inches(1.25), "需要改造", "论文建正常子空间；业务更强调旧型号缺陷库到新型号的共享缺陷子空间。", GOLD)
    card(slide, Inches(8.88), Inches(1.55), Inches(3.7), Inches(1.25), "建议 baseline", "Normal-PCA residual + Defect-prototype projection 双线验证。", BLUE)
    flow(
        slide,
        [
            "旧型号标签\n缺陷原型",
            "新型号无标签\n产线图像",
            "共享投影\n/距离计算",
            "划痕概率\n+热图证据",
        ],
        Inches(0.92),
        Inches(3.38),
        Inches(11.25),
        GOLD
    )
    bullets(
        slide,
        [
            "Normal-PCA：验证新型号正常外观是否能形成稳定残差阈值。",
            "Defect-prototype：验证旧型号划痕特征是否能跨型号投影迁移。",
        ],
        Inches(1.0),
        Inches(5.35),
        Inches(10.8),
        Inches(0.95),
        15
    )
    footer(slide, 11)

    slide = blank(prs)
    title(slide, "实验设计：先回答三个可证伪问题", "先做可解释 baseline，再决定是否引入跨型号域对齐")
    card(slide, Inches(0.78), Inches(1.48), Inches(3.85), Inches(1.18), "实验 1：正常子空间是否成立", "每个新型号取少量正常图拟合 Normal-PCA，测试划痕残差是否稳定高于正常区域。", GREEN)
    card(slide, Inches(4.88), Inches(1.48), Inches(3.85), Inches(1.18), "实验 2：缺陷原型能否迁移", "旧型号划痕 patch 构建 defect prototype / defect PCA，新型号投影后看距离和热图位置。", BLUE)
    card(slide, Inches(8.98), Inches(1.48), Inches(3.85), Inches(1.18), "实验 3：域对齐是否有收益", "旧型号有标签 + 新型号无标签，加入 pull / push 投影目标，对比不对齐 baseline。", GOLD)
    flow(
        slide,
        [
            "Data split\nold/new/type",
            "Feature extraction\npatch/region",
            "Normal residual\nbaseline",
            "Defect projection\nbaseline",
            "UDA ablation\nif needed",
        ],
        Inches(0.8),
        Inches(3.08),
        Inches(11.8),
        BLUE
    )
    bullets(
        slide,
        [
            "评估：I-AUROC、P-AUROC、PRO、误检率、漏检率、阈值稳定性、人工复核成本。",
            "必要可视化：投影前后点云、top patch 位置、异常热图、阈值曲线、跨型号分组结果。",
        ],
        Inches(0.95),
        Inches(5.35),
        Inches(11.3),
        Inches(0.9),
        15
    )
    footer(slide, 12)

    slide = blank(prs)
    title(slide, "风险、限制与下一步问题", "需要把样本条件和验证顺序说清楚")
    card(slide, Inches(0.82), Inches(1.55), Inches(3.65), Inches(1.2), "风险 1：算力", "DINOv2-G 成本高；实际部署可能要用 S/B 模型或蒸馏特征。", RED)
    card(slide, Inches(4.85), Inches(1.55), Inches(3.65), Inches(1.2), "风险 2：线性假设", "PCA 是线性子空间；强反光和复杂纹理可能需要 robust PCA / metric learning。", RED)
    card(slide, Inches(8.88), Inches(1.55), Inches(3.65), Inches(1.2), "风险 3：样本条件", "如果新型号完全无正常样本，Normal-PCA 不成立，需要旧库缺陷子空间或 batched zero-shot。", RED)
    card(slide, Inches(1.0), Inches(3.32), Inches(3.55), Inches(1.3), "问题 1", "业务是否允许每个新型号采少量正常图作为启动样本？", GOLD, 15)
    card(slide, Inches(4.9), Inches(3.32), Inches(3.55), Inches(1.3), "问题 2", "划痕标注应做到 image-level、patch-level，还是只需要少量人工确认点？", GOLD, 15)
    card(slide, Inches(8.8), Inches(3.32), Inches(3.55), Inches(1.3), "问题 3", "第一阶段优先追求可解释 baseline，还是直接做 UDA 共享投影矩阵？", GOLD, 15)
    text_box(slide, "建议：先用 Normal-PCA 做 sanity check，再用旧型号缺陷原型验证跨型号投影；两个 baseline 都站不住时，再进入 UDA 或非线性子空间建模。", Inches(1.05), Inches(5.42), Inches(11.0), Inches(0.85), 17, NAVY, True, PP_ALIGN.CENTER)
    footer(slide, 13)

    prs.save(output_path)
    logger.info("Saved PPTX: %s", output_path)


def parse_args():
    """Parse command-line arguments.

    Parameters:
        None.
    """
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--output",
        default = "2026-04-27/results/subspacead_ppt/SubspaceAD_method_business_v2.pptx",
        help = "Output PPTX path."
    )
    return parser.parse_args()


if __name__ == "__main__":
    logging.basicConfig(
        level = logging.INFO,
        format = '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers = [logging.StreamHandler()]
    )
    args = parse_args()
    build(Path(args.output))
