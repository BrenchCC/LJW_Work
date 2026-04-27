import os
import sys
import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

COLOR_NAVY = RGBColor(15, 23, 42)
COLOR_TEXT = RGBColor(51, 65, 85)
COLOR_MUTED = RGBColor(100, 116, 139)
COLOR_BLUE = RGBColor(30, 64, 175)
COLOR_LIGHT = RGBColor(244, 247, 255)
COLOR_BORDER = RGBColor(226, 232, 240)
COLOR_GOLD = RGBColor(180, 126, 34)
COLOR_WHITE = RGBColor(255, 255, 255)


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size,
    color = COLOR_TEXT,
    bold = False,
    align = PP_ALIGN.LEFT
):
    """Add a text box to one slide.

    Parameters:
        slide: The PowerPoint slide object to update.
        text: Text content to render.
        left: Left position of the text box.
        top: Top position of the text box.
        width: Width of the text box.
        height: Height of the text box.
        font_size: Font size in points.
        color: RGBColor used by the text.
        bold: Whether to render the text in bold.
        align: Paragraph alignment enum.
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = "PingFang SC"
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return shape


def add_bullets(slide, items, left, top, width, height, font_size = 20):
    """Add bullet list text to one slide.

    Parameters:
        slide: The PowerPoint slide object to update.
        items: Sequence of bullet strings.
        left: Left position of the bullet box.
        top: Top position of the bullet box.
        width: Width of the bullet box.
        height: Height of the bullet box.
        font_size: Font size in points.
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.name = "PingFang SC"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = COLOR_TEXT
    return shape


def add_title(slide, title, subtitle = None):
    """Add a consistent title block.

    Parameters:
        slide: The PowerPoint slide object to update.
        title: Main slide title.
        subtitle: Optional small subtitle under the title.
    """
    add_text(slide, title, Inches(0.65), Inches(0.38), Inches(11.9), Inches(0.55), 29, COLOR_NAVY, True)
    line = slide.shapes.add_shape(1, Inches(0.65), Inches(1.05), Inches(11.95), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BLUE
    line.line.color.rgb = COLOR_BLUE
    if subtitle:
        add_text(slide, subtitle, Inches(0.67), Inches(1.13), Inches(11.6), Inches(0.34), 13, COLOR_MUTED)


def add_card(slide, left, top, width, height, title, body, accent = COLOR_BLUE):
    """Add a rounded information card.

    Parameters:
        slide: The PowerPoint slide object to update.
        left: Left position of the card.
        top: Top position of the card.
        width: Width of the card.
        height: Height of the card.
        title: Card title text.
        body: Card body text.
        accent: RGBColor for the accent bar and title.
    """
    card = slide.shapes.add_shape(5, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT
    card.line.color.rgb = COLOR_BORDER
    bar = slide.shapes.add_shape(1, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.color.rgb = accent
    add_text(slide, title, left + Inches(0.18), top + Inches(0.15), width - Inches(0.35), Inches(0.35), 16, accent, True)
    add_text(slide, body, left + Inches(0.18), top + Inches(0.58), width - Inches(0.35), height - Inches(0.75), 14, COLOR_TEXT)


def add_image_fit(slide, image_path, left, top, width, height):
    """Add an image scaled to fit within a target rectangle.

    Parameters:
        slide: The PowerPoint slide object to update.
        image_path: Path to the image file.
        left: Left position of the target rectangle.
        top: Top position of the target rectangle.
        width: Maximum image width.
        height: Maximum image height.
    """
    if not image_path.exists():
        logger.warning("Missing image: %s", image_path)
        return None
    return slide.shapes.add_picture(str(image_path), left, top, width = width, height = height)


def add_footer(slide, number):
    """Add a compact footer to one slide.

    Parameters:
        slide: The PowerPoint slide object to update.
        number: Slide number displayed in the footer.
    """
    add_text(slide, f"SubspaceAD / Feature Subspace Projection  ·  {number:02d}", Inches(0.67), Inches(7.1), Inches(6.0), Inches(0.2), 8, COLOR_MUTED)


def make_slide(prs):
    """Create a blank slide.

    Parameters:
        prs: Presentation object to receive the slide.
    """
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_deck(output_path):
    """Build the offline PowerPoint deck.

    Parameters:
        output_path: Destination .pptx path.
    """
    root = Path(os.getcwd())
    images = root / "2026-04-27" / "results" / "assay_analysis" / "images"
    background = root / "2026-04-27" / "results" / "background_pages"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = make_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_WHITE
    add_text(slide, "SubspaceAD", Inches(0.72), Inches(1.1), Inches(5.8), Inches(0.75), 44, COLOR_NAVY, True)
    add_text(slide, "从正常子空间到特征投影", Inches(0.72), Inches(1.92), Inches(6.6), Inches(0.55), 30, COLOR_BLUE, True)
    add_text(slide, "服务于工业零样本缺陷检测中的特征子空间投影任务", Inches(0.75), Inches(2.7), Inches(6.4), Inches(0.45), 17, COLOR_TEXT)
    add_text(slide, "mentor 论文汇报 / 2026-04-27", Inches(0.75), Inches(6.72), Inches(5.0), Inches(0.28), 11, COLOR_MUTED)
    add_image_fit(slide, background / "page-01.png", Inches(6.25), Inches(0.85), Inches(6.6), Inches(4.0))
    add_card(slide, Inches(7.0), Inches(5.25), Inches(5.35), Inches(1.0), "One-line takeaway", "DINOv2 patch features + PCA normal subspace turns anomaly detection into projection residual scoring.")
    add_footer(slide, 1)

    slide = make_slide(prs)
    add_title(slide, "为什么这篇论文值得看", "从旧型号缺陷库到新产品零样本检测")
    add_card(slide, Inches(0.75), Inches(1.75), Inches(3.7), Inches(1.4), "任务痛点", "新产品无标签，旧缺陷标签直接复用会被材质、光照、背景和纹理变化干扰。")
    add_card(slide, Inches(4.8), Inches(1.75), Inches(3.7), Inches(1.4), "论文价值", "证明强视觉特征配合轻量统计子空间，就能产生稳定的异常分数和定位图。")
    add_card(slide, Inches(8.85), Inches(1.75), Inches(3.7), Inches(1.4), "对本任务启发", "先从图像空间转到特征空间，再用投影距离、残差或分类器做判别。")
    add_image_fit(slide, background / "page-03.png", Inches(1.1), Inches(3.6), Inches(5.25), Inches(2.85))
    add_image_fit(slide, images / "technical_framework.png", Inches(6.75), Inches(3.45), Inches(5.45), Inches(3.05))
    add_footer(slide, 2)

    slide = make_slide(prs)
    add_title(slide, "背景任务：零样本缺陷检测的子空间假设", "物理缺陷相似，成像域发生偏移")
    add_bullets(
        slide,
        [
            "可用资源：旧型号电池已有缺陷样本和标签；新型号只有产线采集图像。",
            "核心假设：划痕的底层物理特征可迁移，但产品外观域需要被剥离或弱化。",
            "子空间路线：预训练 backbone 抽特征，学习共享缺陷子空间，让同类缺陷聚拢、异类样本分离。",
            "在线判别：新样本投影到子空间后，用距离度量或轻量分类器输出缺陷概率。",
        ],
        Inches(0.8),
        Inches(1.65),
        Inches(5.9),
        Inches(4.7),
        18
    )
    add_image_fit(slide, background / "page-09.png", Inches(6.95), Inches(1.65), Inches(5.75), Inches(4.55))
    add_footer(slide, 3)

    slide = make_slide(prs)
    add_title(slide, "SubspaceAD 方法：正常子空间 + 残差评分", "冻结 DINOv2-G，PCA 建模正常 patch 特征")
    add_image_fit(slide, images / "technical_framework.png", Inches(0.8), Inches(1.55), Inches(5.65), Inches(3.2))
    add_image_fit(slide, images / "technical_pca_scoring.png", Inches(6.7), Inches(1.55), Inches(5.65), Inches(3.2))
    add_card(slide, Inches(0.9), Inches(5.15), Inches(3.6), Inches(1.1), "投影", "x_proj = mu + C C^T (x - mu)")
    add_card(slide, Inches(4.9), Inches(5.15), Inches(3.6), Inches(1.1), "评分", "S(x) = ||x - x_proj||^2")
    add_card(slide, Inches(8.9), Inches(5.15), Inches(3.2), Inches(1.1), "解释", "正常 patch 可重构；异常 patch 留下高残差。")
    add_footer(slide, 4)

    slide = make_slide(prs)
    add_title(slide, "实验结论：简单投影能打到强基线", "少样本异常检测与定位结果")
    add_image_fit(slide, images / "table_1_main_comparison.png", Inches(0.75), Inches(1.45), Inches(6.0), Inches(3.5))
    add_image_fit(slide, images / "figure_3_qualitative.png", Inches(6.95), Inches(1.45), Inches(5.65), Inches(3.5))
    add_card(slide, Inches(0.9), Inches(5.25), Inches(2.8), Inches(1.0), "MVTec 1-shot", "I-AUROC 97.1%\nP-AUROC 97.5%")
    add_card(slide, Inches(4.05), Inches(5.25), Inches(2.8), Inches(1.0), "VisA 1-shot", "I-AUROC 93.4%\nP-AUROC 98.2%")
    add_card(slide, Inches(7.2), Inches(5.25), Inches(4.9), Inches(1.0), "谨慎结论", "多数指标达到或超过强基线，但不是每个单项都绝对第一。")
    add_footer(slide, 5)

    slide = make_slide(prs)
    add_title(slide, "关键消融：子空间不是越大越好", "保留完整空间会削弱残差信号")
    add_image_fit(slide, images / "table_4_pca_variance.png", Inches(0.85), Inches(1.55), Inches(3.65), Inches(2.65))
    add_image_fit(slide, images / "figure_4_resolution.png", Inches(4.85), Inches(1.55), Inches(3.65), Inches(2.65))
    add_image_fit(slide, images / "figure_5_backbone.png", Inches(8.85), Inches(1.55), Inches(3.65), Inches(2.65))
    add_card(slide, Inches(0.95), Inches(4.8), Inches(3.45), Inches(1.3), "PCA tau", "tau = 0.95-0.99 较稳；tau = 1.00 时异常也可能被重构，残差失效。")
    add_card(slide, Inches(4.95), Inches(4.8), Inches(3.45), Inches(1.3), "分辨率", "448 px 以上较稳；672 px 是跨 MVTec 与 VisA 的折中。")
    add_card(slide, Inches(8.95), Inches(4.8), Inches(3.45), Inches(1.3), "Backbone", "DINOv2-G 带来上限，小模型可作为速度折中。")
    add_footer(slide, 6)

    slide = make_slide(prs)
    add_title(slide, "和我们的特征子空间投影任务如何对齐", "可借鉴，但不能直接照搬")
    add_card(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(1.25), "一致点", "都从图像空间转到特征空间；都依赖预训练视觉 backbone；都希望用投影距离 / 残差降低标注和训练成本。")
    add_card(slide, Inches(0.8), Inches(3.15), Inches(5.8), Inches(1.25), "差异点", "SubspaceAD 建模正常子空间；我们的背景任务更偏共享缺陷子空间，并且希望新型号零标注。")
    add_card(slide, Inches(0.8), Inches(4.7), Inches(5.8), Inches(1.25), "转化方向", "用 PCA 残差作为 baseline，再扩展到 UDA 聚拢 / 排斥目标和跨型号缺陷原型。")
    add_image_fit(slide, background / "page-04.png", Inches(7.0), Inches(1.6), Inches(5.45), Inches(4.35))
    add_footer(slide, 7)

    slide = make_slide(prs)
    add_title(slide, "下一步：把论文改造成可验证路线", "给 mentor 讨论的实验切入点")
    add_bullets(
        slide,
        [
            "1. 用 DINOv2 / ViT 提取旧型号与新型号 patch 或 region 特征。",
            "2. 建两个 baseline：正常 PCA 残差；缺陷原型子空间投影。",
            "3. 比较距离策略：欧氏、余弦、Mahalanobis、SVM / MLP。",
            "4. 加入域对齐目标：同类缺陷聚拢、背景域排斥、正常残差校准。",
            "5. 可视化检查：投影前后点云、异常热图、阈值敏感性、跨型号泛化。",
        ],
        Inches(0.85),
        Inches(1.55),
        Inches(7.1),
        Inches(4.8),
        17
    )
    add_card(slide, Inches(8.45), Inches(1.75), Inches(3.85), Inches(1.35), "讨论问题", "本阶段优先验证正常子空间残差，还是共享缺陷子空间投影？", COLOR_GOLD)
    add_image_fit(slide, background / "page-11.png", Inches(8.45), Inches(3.5), Inches(3.85), Inches(2.35))
    add_footer(slide, 8)

    prs.save(output_path)
    logger.info("Saved PPTX: %s", output_path)


def parse_args():
    """Parse command-line arguments.

    Parameters:
        None.
    """
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--output",
        default = "2026-04-27/results/subspacead_ppt/SubspaceAD_feature_subspace_report.pptx",
        help = "Output PPTX path."
    )
    return parser.parse_args()


if __name__ == "__main__":
    logging.basicConfig(
        level = logging.INFO,
        format = '%(asctime)s - %(name)s - %(levelname)s - %(message)s',
        handlers = [logging.StreamHandler()]
    )
    args = parse_args()
    build_deck(Path(args.output))
