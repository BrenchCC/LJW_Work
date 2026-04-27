import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Add project root to Python path
sys.path.append(os.getcwd())


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(8, 35, 62)
BLUE = RGBColor(29, 78, 216)
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(71, 85, 105)
LIGHT = RGBColor(244, 247, 251)
LINE = RGBColor(203, 213, 225)
ORANGE = RGBColor(217, 119, 6)
GREEN = RGBColor(15, 118, 110)
RED = RGBColor(185, 28, 28)
WHITE = RGBColor(255, 255, 255)


def set_text_frame(text_frame, font_size = 16, color = INK, bold = False):
    """Set font style for every paragraph and run in a text frame.

    Parameters:
        text_frame: The python-pptx text frame to style.
        font_size: Font size in points.
        color: RGBColor used for text.
        bold: Whether text should use bold weight.
    """
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_textbox(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size = 16,
    color = INK,
    bold = False,
    align = PP_ALIGN.LEFT
):
    """Add a text box with consistent font styling.

    Parameters:
        slide: The slide object receiving the text box.
        text: Text content to place in the box.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        font_size: Font size in points.
        color: RGBColor used for text.
        bold: Whether text should use bold weight.
        align: Paragraph alignment from python-pptx.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.02)
    text_frame.margin_right = Inches(0.02)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    set_text_frame(text_frame, font_size = font_size, color = color, bold = bold)
    return shape


def add_header(slide, section, title, index, total):
    """Add the recurring slide header and footer.

    Parameters:
        slide: The slide object receiving the chrome.
        section: Short section label.
        title: Main slide title.
        index: Current slide index, starting from 1.
        total: Total slide count.
    """
    add_textbox(slide, section.upper(), 0.55, 0.34, 4.2, 0.25, font_size = 8, color = BLUE, bold = True)
    add_textbox(slide, title, 0.55, 0.78, 11.6, 0.64, font_size = 27, color = NAVY, bold = True)
    add_textbox(
        slide,
        f"{index} / {total}",
        12.15,
        7.04,
        0.7,
        0.2,
        font_size = 8,
        color = MUTED,
        align = PP_ALIGN.RIGHT
    )
    line = slide.shapes.add_shape(1, Inches(0), Inches(7.28), Inches(12.0 * index / total), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.color.rgb = NAVY


def add_panel(slide, x, y, w, h, fill = WHITE, line = LINE):
    """Add a rounded content panel.

    Parameters:
        slide: The slide object receiving the panel.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        fill: RGBColor used as panel fill.
        line: RGBColor used as panel border.
    """
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_bullets(slide, items, x, y, w, h, font_size = 14, color = MUTED):
    """Add a bullet list with consistent spacing.

    Parameters:
        slide: The slide object receiving the bullet list.
        items: Bullet strings.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        font_size: Font size in points.
        color: RGBColor used for bullet text.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.02)
    for idx, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(5)
    return shape


def add_image_fit(slide, image_path, x, y, w, h):
    """Add an image without cropping or stretching.

    Parameters:
        slide: The slide object receiving the image.
        image_path: Path to the source image.
        x: Left boundary in inches.
        y: Top boundary in inches.
        w: Maximum width in inches.
        h: Maximum height in inches.
    """
    image = Image.open(image_path)
    img_w, img_h = image.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
    else:
        draw_h = h
        draw_w = h * img_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(draw_x),
        Inches(draw_y),
        width = Inches(draw_w),
        height = Inches(draw_h)
    )


def add_card(slide, title, body, x, y, w, h, accent = BLUE):
    """Add a titled card.

    Parameters:
        slide: The slide object receiving the card.
        title: Card heading.
        body: Card body text.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        accent: RGBColor used for the title.
    """
    add_panel(slide, x, y, w, h)
    add_textbox(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.34, font_size = 14, color = accent, bold = True)
    add_textbox(slide, body, x + 0.18, y + 0.68, w - 0.36, h - 0.8, font_size = 12, color = MUTED)


def add_table(slide, rows, x, y, w, h, col_widths = None):
    """Add a simple editable table.

    Parameters:
        slide: The slide object receiving the table.
        rows: Two-dimensional list of cell strings.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        col_widths: Optional relative column widths.
    """
    table_shape = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )
    table = table_shape.table
    if col_widths:
        total_width = sum(col_widths)
        for idx, rel_width in enumerate(col_widths):
            table.columns[idx].width = Inches(w * rel_width / total_width)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(9 if row_idx == 0 else 8)
                paragraph.font.bold = row_idx == 0
                paragraph.font.color.rgb = WHITE if row_idx == 0 else INK
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row_idx == 0 else WHITE
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return table_shape


def build_deck(output_path):
    """Build the rich SubspaceAD project report deck.

    Parameters:
        output_path: Path where the PowerPoint file should be saved.
    """
    base = Path("2026-04-27/results/assay_analysis/images")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 15

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, "PAPER REPORT · 2026-04-27", 0.55, 0.45, 5.2, 0.28, font_size = 9, color = BLUE, bold = True)
    add_textbox(slide, "SubspaceAD：\n正常子空间建模与特征投影实验", 0.55, 1.08, 6.65, 1.85, font_size = 30, color = NAVY, bold = True)
    add_textbox(
        slide,
        "Training-Free Few-Shot Anomaly Detection via Subspace Modeling\n目标：说清方法，明确可验证实验，判断它是否适合电池划痕少样本检测。",
        0.58,
        3.18,
        5.9,
        0.95,
        font_size = 15,
        color = MUTED
    )
    add_panel(slide, 7.0, 1.15, 5.6, 2.3)
    add_image_fit(slide, base / "header.png", 7.18, 1.3, 5.25, 1.92)
    for idx, label in enumerate(["DINOv2 patch features", "PCA normal subspace", "Residual anomaly score"]):
        add_card(slide, f"0{idx + 1}", label, 0.62 + idx * 2.05, 4.72, 1.85, 1.08, accent = BLUE)
    add_textbox(slide, "核心问题：强基础模型特征足够好时，异常检测是否还需要复杂训练和大型 memory bank？", 0.7, 6.22, 11.8, 0.35, font_size = 15, color = NAVY, bold = True)
    add_header(slide, "", "", 1, total)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_header(slide, "01 · motivation", "论文要解决什么问题", 2, total)
    add_card(slide, "工业异常检测的痛点", "正常样本相对容易拿到；缺陷样本稀少、类型长尾，且细微划痕、污染、错位等很难靠文本提示完整描述。", 0.7, 1.8, 3.8, 1.65, accent = RED)
    add_card(slide, "现有少样本方法的复杂度", "PatchCore / SPADE 依赖 memory bank；WinCLIP / PromptAD / IIPAD 依赖文本提示或调优；多阶段流程增加部署和维护成本。", 4.8, 1.8, 3.8, 1.65, accent = ORANGE)
    add_card(slide, "SubspaceAD 的回答", "如果 DINOv2 的 dense patch feature 足够强，用 PCA 描述正常变化，再用重构残差即可形成异常分数。", 8.9, 1.8, 3.7, 1.65, accent = GREEN)
    add_panel(slide, 0.7, 4.12, 11.9, 1.35, fill = LIGHT)
    add_textbox(slide, "报告视角", 0.95, 4.38, 1.2, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "这不是单纯复述 SOTA，而是把它拆成可落地模块：特征抽取、子空间拟合、残差评分、图像级聚合、阈值校准。", 2.1, 4.38, 9.9, 0.55, font_size = 15, color = INK)
    add_bullets(slide, ["项目对应：新型号电池缺陷标签少，先验证特征空间中的投影距离是否比图像空间更稳定。", "实验优先级：先做能解释失败原因的 baseline，再决定是否上 UDA 或更复杂模型。"], 0.85, 5.82, 11.4, 0.8, font_size = 13)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_header(slide, "02 · positioning", "方法定位：从 memory bank 到统计子空间", 3, total)
    rows = [
        ["路线", "代表方法", "核心机制", "部署风险"],
        ["Memory bank", "SPADE / PatchCore / AnomalyDINO", "保存正常 patch 特征，最近邻匹配", "存储与检索开销随样本增长"],
        ["VLM prompt", "WinCLIP / PromptAD / IIPAD", "用文本或实例提示对齐异常语义", "依赖 prompt、辅助数据或调优"],
        ["SubspaceAD", "DINOv2 + PCA", "拟合正常低维子空间，残差即异常", "主要瓶颈转为 backbone 特征质量"]
    ]
    add_table(slide, rows, 0.75, 1.75, 11.8, 2.35, col_widths = [1.15, 2.1, 2.35, 2.1])
    add_panel(slide, 0.75, 4.65, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "方法主张", 1.05, 4.92, 1.2, 0.32, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "把“正常性”建模从样本检索改成分布级 PCA：模型只保留均值 μ 与主方向 C，异常是无法被 C 充分重构的 patch。", 2.25, 4.92, 9.55, 0.5, font_size = 15, color = INK)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_header(slide, "03 · pipeline", "整体流程：两个阶段，一个残差图", 4, total)
    add_panel(slide, 0.75, 1.62, 7.15, 3.2)
    add_image_fit(slide, base / "technical_framework.png", 0.95, 1.82, 6.75, 2.8)
    add_card(slide, "Fitting", "k 张正常图 + 旋转增强 → DINOv2-G 中间层 patch feature → PCA 拟合正常子空间。", 8.25, 1.62, 4.25, 1.25, accent = BLUE)
    add_card(slide, "Inference", "测试图走同一个 feature extractor，投影到正常子空间；每个 patch 的重构误差形成 anomaly map。", 8.25, 3.1, 4.25, 1.25, accent = GREEN)
    add_panel(slide, 0.75, 5.35, 11.75, 0.75, fill = LIGHT)
    add_textbox(slide, "一句话", 1.0, 5.52, 0.9, 0.3, font_size = 13, color = NAVY, bold = True)
    add_textbox(slide, "SubspaceAD 不学习“异常长什么样”，而是学习“正常变化能张成什么空间”。", 1.9, 5.5, 9.8, 0.35, font_size = 15, color = INK)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_header(slide, "04 · feature extraction", "特征抽取：为什么用 DINOv2 中间层", 5, total)
    add_card(slide, "Backbone", "冻结 DINOv2-G/14，输入分辨率实验最终采用 672 px。PCA 本身很轻，主要计算来自 backbone forward。", 0.75, 1.72, 3.8, 1.45, accent = BLUE)
    add_card(slide, "Layer aggregation", "不是只用最后一层，而是 mean-pool 第 22-28 层。中间层保留局部结构与语义，适合 patch-level 缺陷。", 4.78, 1.72, 3.8, 1.45, accent = GREEN)
    add_card(slide, "Augmentation", "每张正常图生成 Na = 30 个旋转视图，构建更稳定的正常协方差估计。", 8.82, 1.72, 3.65, 1.45, accent = ORANGE)
    rows = [
        ["Aggregation", "MVTec I-AUROC", "MVTec PRO", "VisA I-AUROC", "VisA PRO"],
        ["Mean-pool Middle-7", "97.9", "94.3", "94.7", "93.8"],
        ["Mean-pool Final-7", "98.2", "92.3", "91.9", "90.3"],
        ["Concat Middle-7", "98.6", "94.2", "93.8", "93.1"],
        ["Last layer only", "97.5", "91.4", "89.1", "88.3"]
    ]
    add_table(slide, rows, 1.0, 3.75, 11.1, 1.8, col_widths = [2.3, 1.35, 1.15, 1.35, 1.15])
    add_textbox(slide, "解释：最后层偏类别语义，容易损失细节；中间层能更好保留划痕、边界、局部纹理这类工业缺陷信号。", 1.0, 6.0, 11.1, 0.42, font_size = 14, color = NAVY, bold = True)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_header(slide, "05 · PCA model", "PCA 正常子空间：模型参数很少", 6, total)
    add_panel(slide, 0.75, 1.65, 5.8, 3.85, fill = LIGHT)
    add_textbox(slide, "Normal feature model", 1.05, 1.95, 3.0, 0.3, font_size = 16, color = NAVY, bold = True)
    add_textbox(slide, "x = μ + C z + ε\nz ~ N(0, I_r)\nε ~ N(0, σ² I)", 1.05, 2.55, 4.7, 1.0, font_size = 24, color = INK, bold = True)
    add_textbox(slide, "C ∈ R^(D×r) 是前 r 个主方向，r 由解释方差阈值 τ 决定。论文默认 τ = 0.99。", 1.05, 4.12, 4.75, 0.65, font_size = 14, color = MUTED)
    add_card(slide, "保存什么", "每个类别只需要 μ 和 C；论文估计通常小于 1 MB。", 7.0, 1.72, 2.55, 1.25, accent = BLUE)
    add_card(slide, "为什么可解释", "异常不是一个黑盒分类结果，而是特征在正常主子空间之外的能量。", 9.9, 1.72, 2.55, 1.25, accent = GREEN)
    add_card(slide, "为什么适合作 baseline", "正常样本少时依然能快速拟合，且可以按类别、产线、型号独立管理。", 7.0, 3.32, 5.45, 1.35, accent = ORANGE)
    add_panel(slide, 7.0, 5.2, 5.45, 0.9)
    add_textbox(slide, "项目含义", 7.25, 5.42, 1.0, 0.3, font_size = 13, color = NAVY, bold = True)
    add_textbox(slide, "先定义“新型号正常外观的可接受变化”，再把偏离该空间的 patch 作为疑似缺陷。", 8.2, 5.38, 3.95, 0.42, font_size = 13, color = INK)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_header(slide, "06 · scoring", "推理评分：投影残差、TVaR 与像素定位", 7, total)
    add_panel(slide, 0.75, 1.65, 5.35, 4.35)
    add_image_fit(slide, base / "technical_pca_scoring.png", 0.98, 1.82, 4.85, 4.0)
    add_card(slide, "Patch-level", "x_proj = μ + C C^T (x_p - μ)\nS(x_p) = ||x_p - x_proj||²", 6.55, 1.68, 5.8, 1.15, accent = BLUE)
    add_card(slide, "Image-level", "用 top ρ% patch 分数的均值作为图像级分数；论文 ρ = 1%，对稀疏划痕更敏感。", 6.55, 3.08, 5.8, 1.15, accent = GREEN)
    add_card(slide, "Pixel-level", "patch anomaly map 双线性上采样到原图分辨率，再用 σ = 4 的 Gaussian filter 平滑。", 6.55, 4.48, 5.8, 1.15, accent = ORANGE)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_header(slide, "07 · experiments", "实验设置：少样本检测与定位", 8, total)
    add_card(slide, "Datasets", "MVTec-AD：15 类工业对象；VisA：高分辨率、更复杂背景和异常类型。", 0.75, 1.7, 3.75, 1.35, accent = BLUE)
    add_card(slide, "Shots", "1-shot / 2-shot / 4-shot；每个设置 5 个随机种子，报告均值和标准差。", 4.82, 1.7, 3.75, 1.35, accent = GREEN)
    add_card(slide, "Metrics", "图像级 AUROC / AUPR；像素级 AUROC / PRO。原始分数用于评估，归一化图用于可视化。", 8.9, 1.7, 3.65, 1.35, accent = ORANGE)
    add_panel(slide, 0.75, 3.65, 11.8, 1.75, fill = LIGHT)
    add_textbox(slide, "对比对象", 1.05, 3.95, 1.3, 0.3, font_size = 15, color = NAVY, bold = True)
    add_bullets(
        slide,
        [
            "Memory bank / retrieval：SPADE、PatchCore、AnomalyDINO",
            "Reconstruction：FastRecon",
            "Vision-language：WinCLIP、PromptAD、IIPAD"
        ],
        2.35,
        3.92,
        9.3,
        0.9,
        font_size = 13,
        color = INK
    )
    add_textbox(slide, "汇报口径：强调方法极简但效果强，同时保留“不是每个单项都绝对第一”的谨慎表述。", 1.0, 6.08, 11.2, 0.36, font_size = 14, color = NAVY, bold = True)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_header(slide, "08 · main results", "主结果：简单投影能达到强基线", 9, total)
    add_panel(slide, 0.65, 1.55, 4.05, 4.7)
    metrics = [
        ("97.1%", "MVTec-AD 1-shot image AUROC"),
        ("97.5%", "MVTec-AD 1-shot pixel AUROC"),
        ("93.4%", "VisA 1-shot image AUROC"),
        ("98.2%", "VisA 1-shot pixel AUROC")
    ]
    for idx, metric in enumerate(metrics):
        add_textbox(slide, metric[0], 0.95, 1.85 + idx * 1.0, 1.3, 0.4, font_size = 22, color = NAVY, bold = True)
        add_textbox(slide, metric[1], 2.25, 1.92 + idx * 1.0, 2.05, 0.32, font_size = 10.5, color = MUTED)
    add_panel(slide, 5.0, 1.55, 7.55, 4.7)
    add_image_fit(slide, base / "table_1_main_comparison.png", 5.2, 1.72, 7.15, 4.35)
    add_textbox(slide, "结论：SubspaceAD 在多数少样本指标上领先或持平；VisA 1-shot 图像级 AUROC 相比 AnomalyDINO 提升约 6 个百分点。", 0.9, 6.43, 11.4, 0.36, font_size = 13, color = NAVY, bold = True)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_header(slide, "09 · localization", "定性定位：残差图不是只做图像级分数", 10, total)
    add_panel(slide, 0.7, 1.58, 11.9, 3.95)
    add_image_fit(slide, base / "figure_3_qualitative.png", 0.95, 1.8, 11.4, 3.5)
    add_card(slide, "观察 1", "在 candle / capsules / cashew 等 VisA 示例中，背景误激活更少，缺陷区域更集中。", 0.75, 5.85, 3.7, 0.75, accent = BLUE)
    add_card(slide, "观察 2", "在 bottle / hazelnut / screw 等 MVTec 示例中，残差热图与 GT 缺陷区域对齐更好。", 4.82, 5.85, 3.7, 0.75, accent = GREEN)
    add_card(slide, "可验证点", "如果热图稳定，后续可接 ROI 级复检、缺陷面积统计和人工复核优先级排序。", 8.9, 5.85, 3.65, 0.75, accent = ORANGE)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_header(slide, "10 · ablations", "消融：什么因素真正决定效果", 11, total)
    add_panel(slide, 0.72, 1.55, 3.65, 3.1)
    add_image_fit(slide, base / "figure_4_resolution.png", 0.92, 1.72, 3.25, 2.55)
    add_textbox(slide, "分辨率：VisA 从 256 到 448/672 明显改善；MVTec 在 448 以上较稳。", 0.9, 4.25, 3.25, 0.35, font_size = 10.5, color = MUTED)
    add_panel(slide, 4.85, 1.55, 3.65, 3.1)
    add_image_fit(slide, base / "figure_5_backbone.png", 5.05, 1.72, 3.25, 2.55)
    add_textbox(slide, "backbone：DINOv2-G 整体优于 S/B/L，小模型可做部署折中。", 5.02, 4.25, 3.25, 0.35, font_size = 10.5, color = MUTED)
    add_panel(slide, 8.98, 1.55, 3.45, 3.1)
    add_image_fit(slide, base / "table_4_pca_variance.png", 9.16, 1.72, 3.1, 2.55)
    add_textbox(slide, "PCA 阈值：τ = 1.00 会把异常也重构掉，残差判别显著变弱。", 9.14, 4.25, 3.05, 0.35, font_size = 10.5, color = MUTED)
    add_panel(slide, 0.75, 5.42, 11.7, 0.9, fill = LIGHT)
    add_textbox(slide, "关键理解", 1.0, 5.65, 1.25, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "性能不是来自 PCA 本身，而来自“强特征 + 合适残差空间”。子空间必须压缩正常主变化，同时保留暴露异常的正交方向。", 2.25, 5.6, 9.6, 0.42, font_size = 13.5, color = INK)

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_header(slide, "11 · efficiency & limits", "复杂度与边界：适合做一阶段 baseline", 12, total)
    add_card(slide, "Storage", "每类模型只保存 μ 与 C，通常 < 1 MB；相比 patch memory bank 更易版本管理。", 0.75, 1.75, 3.7, 1.4, accent = BLUE)
    add_card(slide, "Latency", "672×672 图像约 300 ms / img；其中 DINOv2-G forward 约 226 ms，投影和评分约 74 ms。", 4.82, 1.75, 3.7, 1.4, accent = GREEN)
    add_card(slide, "Scoring head", "H100 上 SubspaceAD scoring 约 74.1 ms / img，且对 1-4 shot 基本不随 K 增长。", 8.9, 1.75, 3.65, 1.4, accent = ORANGE)
    add_panel(slide, 0.75, 3.75, 11.8, 1.85, fill = LIGHT)
    add_textbox(slide, "论文暴露的限制", 1.0, 4.02, 1.55, 0.3, font_size = 14, color = RED, bold = True)
    add_bullets(
        slide,
        [
            "逻辑 / 结构异常较难：例如 Transistor 的 missing / misplaced component，patch 外观方法不显式建模结构关系。",
            "正常样本类内变化大或背景伪影多时，少量样本很难形成紧致正常子空间。",
            "旋转增强不是无脑使用；如果方向本身定义正常状态，旋转会把方向异常学成正常变化。"
        ],
        2.55,
        3.95,
        9.3,
        1.05,
        font_size = 12.2,
        color = INK
    )
    add_textbox(slide, "落地判断：先用于局部外观类缺陷 baseline；涉及装配关系、方向、缺件时，需要加入几何/结构先验。", 0.95, 6.05, 11.2, 0.36, font_size = 13.5, color = NAVY, bold = True)

    # Slide 13
    slide = prs.slides.add_slide(blank)
    add_header(slide, "12 · mapping to our task", "映射到电池划痕任务：正常子空间 vs 共享缺陷子空间", 13, total)
    rows = [
        ["维度", "SubspaceAD", "电池划痕任务可用改造"],
        ["训练数据", "k 张新类别正常图", "新型号少量正常图 + 旧型号划痕标注"],
        ["子空间含义", "正常外观主变化", "正常子空间 + 共享划痕原型子空间"],
        ["异常分数", "到正常子空间的残差", "残差、到划痕原型距离、轻量分类器联合"],
        ["泛化挑战", "类内变化与背景干扰", "型号域差异、材质纹理、光照与工位变化"]
    ]
    add_table(slide, rows, 0.75, 1.65, 11.8, 2.6, col_widths = [1.15, 2.25, 3.2])
    add_panel(slide, 0.75, 4.75, 11.8, 1.25, fill = LIGHT)
    add_textbox(slide, "下一步判断", 1.0, 5.05, 1.3, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "不要一开始就押注复杂 UDA；先并行做两个 baseline：正常 PCA 残差与旧型号划痕原型投影。二者能回答“新型号无标签时，哪种投影信号更稳定”。", 2.25, 4.95, 9.7, 0.56, font_size = 13.5, color = INK)

    # Slide 14
    slide = prs.slides.add_slide(blank)
    add_header(slide, "13 · experiment workflow", "实际验证流程：从图像到误差分析", 14, total)
    flow = [
        ("1", "采样", "新型号：少量正常图 + 未标注混合图；旧型号：保留划痕标注样本。"),
        ("2", "特征", "DINOv2/ViT 提取 patch 或 ROI 特征，保存图像坐标映射，便于回看热图。"),
        ("3", "建模", "分别拟合正常 PCA、旧型号划痕原型、共享缺陷子空间。"),
        ("4", "评分", "输出 residual score、prototype distance、top-ρ patch score。"),
        ("5", "评估", "用人工复核或少量标签计算 top-k precision、recall、误报来源。"),
        ("6", "复盘", "按失败类型决定：调特征层、调阈值、加 ROI、还是引入域对齐。")
    ]
    for idx, item in enumerate(flow):
        x = 0.75 + (idx % 3) * 4.05
        y = 1.75 + (idx // 3) * 2.0
        add_panel(slide, x, y, 3.55, 1.45)
        add_textbox(slide, item[0], x + 0.18, y + 0.18, 0.42, 0.34, font_size = 14, color = WHITE, bold = True, align = PP_ALIGN.CENTER)
        badge = slide.shapes.add_shape(5, Inches(x + 0.12), Inches(y + 0.14), Inches(0.45), Inches(0.36))
        badge.fill.solid()
        badge.fill.fore_color.rgb = NAVY
        badge.line.color.rgb = NAVY
        add_textbox(slide, item[1], x + 0.72, y + 0.2, 2.3, 0.3, font_size = 15, color = NAVY, bold = True)
        add_textbox(slide, item[2], x + 0.2, y + 0.72, 3.05, 0.48, font_size = 11.2, color = MUTED)
    add_textbox(slide, "原则：先离线验证，不急着承诺上线效果；每个失败样本都要能归因到特征、子空间、阈值或数据覆盖问题。", 0.85, 6.3, 11.4, 0.34, font_size = 13.5, color = NAVY, bold = True)

    # Slide 15
    slide = prs.slides.add_slide(blank)
    add_header(slide, "14 · experiment ideas", "实验想法与需要回答的问题", 15, total)
    rows = [
        ["实验", "要验证什么", "做法", "看什么结果"],
        ["E1 正常 PCA", "新型号少量正常图是否足够定义正常空间", "k = 1/2/4/8；比较 τ = 0.95/0.99/1.00", "热图是否覆盖划痕；误报是否来自背景"],
        ["E2 缺陷原型", "旧型号划痕特征能否迁移到新型号", "旧型号标注 patch 建原型/子空间", "新型号 top-k 疑似样本命中率"],
        ["E3 特征选择", "层、backbone、分辨率哪个最影响结果", "DINOv2-S/B/L/G；middle/final layer；448/672", "精度、速度、显存折中"],
        ["E4 阈值策略", "图像级分数如何稳定", "top-ρ TVaR；ρ = 0.5/1/2/5%", "召回优先下的误报数量"],
        ["E5 失败归因", "失败主要来自特征还是数据域", "UMAP/PCA 点云 + 误报 patch 回看", "是否需要 ROI、域对齐或结构先验"]
    ]
    add_table(slide, rows, 0.7, 1.55, 11.95, 3.2, col_widths = [1.2, 2.0, 2.2, 2.3])
    add_panel(slide, 0.75, 5.25, 11.85, 1.1, fill = LIGHT)
    add_textbox(slide, "需要先回答的问题", 1.0, 5.55, 2.1, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "正常子空间残差能不能单独解决？旧型号划痕是否真的可迁移？需要多少正常图才稳定？先做 ROI 还是整图？阈值按召回还是按复核量定？", 3.05, 5.43, 8.9, 0.56, font_size = 13.5, color = INK)

    prs.save(output_path)


if __name__ == "__main__":
    build_deck("2026-04-27/results/html_ppt/SubspaceAD_project_experiment_report.pptx")
