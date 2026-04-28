import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.util import Inches
from pptx.dml.color import RGBColor

sys.path.append(os.getcwd())


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(8, 35, 62)
BLUE = RGBColor(29, 78, 216)
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(71, 85, 105)
LIGHT = RGBColor(244, 247, 251)
LINE = RGBColor(203, 213, 225)
ORANGE = RGBColor(217, 119, 6)
GREEN = RGBColor(15, 118, 110)
RED = RGBColor(185, 28, 28)
WHITE = RGBColor(255, 255, 255)


def set_text_frame(text_frame, font_size = 16, color = INK, bold = False):
    """Set font style for every paragraph and run in a text frame.

    Parameters:
        text_frame: The python-pptx text frame to style.
        font_size: Font size in points.
        color: RGBColor used for text.
        bold: Whether text should use bold weight.
    """
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        for run in paragraph.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold


def add_textbox(
    slide,
    text,
    x,
    y,
    w,
    h,
    font_size = 16,
    color = INK,
    bold = False,
    align = PP_ALIGN.LEFT
):
    """Add a text box with consistent font styling.

    Parameters:
        slide: The slide object receiving the text box.
        text: Text content to place in the box.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        font_size: Font size in points.
        color: RGBColor used for text.
        bold: Whether text should use bold weight.
        align: Paragraph alignment from python-pptx.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.02)
    text_frame.margin_right = Inches(0.02)
    text_frame.margin_top = Inches(0.02)
    text_frame.margin_bottom = Inches(0.02)
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    set_text_frame(text_frame, font_size = font_size, color = color, bold = bold)
    return shape


def add_header(slide, section, title, index, total):
    """Add the recurring slide header and footer.

    Parameters:
        slide: The slide object receiving the chrome.
        section: Short section label.
        title: Main slide title.
        index: Current slide index, starting from 1.
        total: Total slide count.
    """
    add_textbox(slide, section.upper(), 0.55, 0.34, 4.2, 0.25, font_size = 8, color = BLUE, bold = True)
    add_textbox(slide, title, 0.55, 0.78, 11.6, 0.64, font_size = 27, color = NAVY, bold = True)
    add_textbox(
        slide,
        f"{index} / {total}",
        12.15,
        7.04,
        0.7,
        0.2,
        font_size = 8,
        color = MUTED,
        align = PP_ALIGN.RIGHT
    )
    line = slide.shapes.add_shape(1, Inches(0), Inches(7.28), Inches(12.0 * index / total), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.color.rgb = NAVY


def add_panel(slide, x, y, w, h, fill = WHITE, line = LINE):
    """Add a rounded content panel.

    Parameters:
        slide: The slide object receiving the panel.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        fill: RGBColor used as panel fill.
        line: RGBColor used as panel border.
    """
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_bullets(slide, items, x, y, w, h, font_size = 14, color = MUTED):
    """Add a bullet list with consistent spacing.

    Parameters:
        slide: The slide object receiving the bullet list.
        items: Bullet strings.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        font_size: Font size in points.
        color: RGBColor used for bullet text.
    """
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.02)
    for idx, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(5)
    return shape


def add_image_fit(slide, image_path, x, y, w, h):
    """Add an image without cropping or stretching.

    Parameters:
        slide: The slide object receiving the image.
        image_path: Path to the source image.
        x: Left boundary in inches.
        y: Top boundary in inches.
        w: Maximum width in inches.
        h: Maximum height in inches.
    """
    image = Image.open(image_path)
    img_w, img_h = image.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
    else:
        draw_h = h
        draw_w = h * img_ratio
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    return slide.shapes.add_picture(
        str(image_path),
        Inches(draw_x),
        Inches(draw_y),
        width = Inches(draw_w),
        height = Inches(draw_h)
    )


def add_card(slide, title, body, x, y, w, h, accent = BLUE):
    """Add a titled card.

    Parameters:
        slide: The slide object receiving the card.
        title: Card heading.
        body: Card body text.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        accent: RGBColor used for the title.
    """
    add_panel(slide, x, y, w, h)
    add_textbox(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.34, font_size = 14, color = accent, bold = True)
    add_textbox(slide, body, x + 0.18, y + 0.68, w - 0.36, h - 0.8, font_size = 12, color = MUTED)


def add_table(slide, rows, x, y, w, h, col_widths = None):
    """Add a simple editable table.

    Parameters:
        slide: The slide object receiving the table.
        rows: Two-dimensional list of cell strings.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        col_widths: Optional relative column widths.
    """
    table_shape = slide.shapes.add_table(
        len(rows),
        len(rows[0]),
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h)
    )
    table = table_shape.table
    if col_widths:
        total_width = sum(col_widths)
        for idx, rel_width in enumerate(col_widths):
            table.columns[idx].width = Inches(w * rel_width / total_width)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(9 if row_idx == 0 else 8)
                paragraph.font.bold = row_idx == 0
                paragraph.font.color.rgb = WHITE if row_idx == 0 else INK
                paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if row_idx == 0 else WHITE
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return table_shape


def add_badge(slide, text, x, y, w, h, bg_color = NAVY):
    """Add a colored badge with text.

    Parameters:
        slide: The slide object receiving the badge.
        text: Badge text.
        x: Left position in inches.
        y: Top position in inches.
        w: Width in inches.
        h: Height in inches.
        bg_color: RGBColor used as badge fill.
    """
    badge = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.color.rgb = bg_color
    add_textbox(slide, text, x + 0.02, y + 0.02, w - 0.04, h - 0.04, font_size = 14, color = WHITE, bold = True, align = PP_ALIGN.CENTER)


def build_deck(output_path):
    """Build the ADDA paper report deck.

    Parameters:
        output_path: Path where the PowerPoint file should be saved.
    """
    base = Path("2026-04-28/results/1702-05464v1/images")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 18

    # ── Slide 1: Cover ──
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, "MENTOR PAPER REPORT · 2026-04-28", 0.55, 0.45, 5.2, 0.28, font_size = 9, color = BLUE, bold = True)
    add_textbox(slide, "ADDA：\n对抗判别式域适应", 0.55, 1.08, 6.65, 1.85, font_size = 30, color = NAVY, bold = True)
    add_textbox(
        slide,
        "Adversarial Discriminative Domain Adaptation\n目标：说清域适应怎么做、为什么有效、对电池缺陷检测能不能用。",
        0.58,
        3.18,
        5.9,
        0.95,
        font_size = 15,
        color = MUTED
    )
    add_panel(slide, 7.0, 1.15, 5.6, 2.3)
    add_image_fit(slide, base / "header.png", 7.18, 1.3, 5.25, 1.92)
    for idx, label in enumerate(["无监督域适应", "对抗学习", "特征子空间对齐"]):
        add_card(slide, f"0{idx + 1}", label, 0.62 + idx * 2.05, 4.72, 1.85, 1.08, accent = BLUE)
    add_textbox(slide, "核心问题：源域有标签、目标域无标签时，如何让分类器跨域迁移？", 0.7, 6.22, 11.8, 0.35, font_size = 15, color = NAVY, bold = True)
    add_header(slide, "", "", 1, total)

    # ── Slide 2: Outline ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "overview", "今天聊什么", 2, total)
    sections = [
        ("1", "问题与动机", "域偏移为什么让模型失效；我们的缺陷检测痛点。", NAVY),
        ("2", "ADDA 方法", "三阶段训练流程、对抗对齐原理、权重不共享的设计。", BLUE),
        ("3", "实验证据", "数字数据集跨域结果、NYUD 跨模态迁移、对比分析。", GREEN),
        ("4", "应用讨论", "能否用到电池缺陷检测、局限与风险、组合方案建议。", ORANGE)
    ]
    for idx, (num, title, body, accent) in enumerate(sections):
        x = 0.75 + idx * 3.05
        add_panel(slide, x, 1.85, 2.75, 3.8)
        add_badge(slide, num, x + 0.18, 2.1, 0.5, 0.4, bg_color = accent)
        add_textbox(slide, title, x + 0.18, 2.75, 2.35, 0.4, font_size = 18, color = NAVY, bold = True)
        add_textbox(slide, body, x + 0.18, 3.3, 2.35, 1.8, font_size = 13, color = MUTED)
    add_panel(slide, 0.75, 6.1, 11.8, 0.75, fill = LIGHT)
    add_textbox(slide, "核心结论先剧透", 1.0, 6.28, 1.8, 0.3, font_size = 13, color = NAVY, bold = True)
    add_textbox(slide, "ADDA 适合做域适应基线，但不能单独解决“第六款没有划痕样本”的根本风险。", 2.8, 6.28, 9.2, 0.35, font_size = 14, color = INK)

    # ── Slide 3: Our Problem ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "01 · the problem", "我们的痛点：跨产品缺陷泛化差", 3, total)
    add_card(slide, "前五款电池", "有缺陷样本与标签。划痕、正常、其他缺陷类别提供监督信号。分类器在五款上表现很好。", 0.75, 1.75, 3.7, 1.65, accent = BLUE)
    add_card(slide, "直接迁移", "泛化性明显下降。第六款电池的材质、光照、纹理变了，特征空间偏移，分类器失准。", 4.8, 1.75, 3.7, 1.65, accent = ORANGE)
    add_card(slide, "第六款电池", "没有划痕标签。暂时无法重新标注，需要无监督或半监督方法弥合域差距。", 8.85, 1.75, 3.7, 1.65, accent = RED)
    add_panel(slide, 0.75, 3.85, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "本质问题", 1.0, 4.12, 1.2, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "产品外观域特征子空间压过了缺陷语义子空间；分类器学到的是“哪款产品”，而不是“有什么缺陷”。", 2.2, 4.08, 9.6, 0.55, font_size = 15, color = INK)
    add_bullets(slide, [
        "域适应要做的事情：把目标域特征簇拉到源域特征簇附近，让“划痕”和“正常”的判别边界在目标域上也能生效。",
        "ADDA 提供了一个清晰的对齐思路：用对抗训练拉近两个域的特征分布。"
    ], 0.85, 5.7, 11.4, 0.8, font_size = 13)

    # ── Slide 4: Domain Shift ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "02 · domain shift", "什么是域偏移：同样的划痕，不同的“长相”", 4, total)
    add_card(slide, "源域：前五款电池", "相同缺陷在不同产品上成像差异大\n光照、材质、纹理、相机参数不同\n特征空间按产品聚类", 0.75, 1.75, 5.2, 1.8, accent = BLUE)
    add_card(slide, "目标域：第六款电池", "划痕物理形态类似\n但图像特征分布与前五款差异大\n源域分类器无法覆盖", 7.35, 1.75, 5.2, 1.8, accent = RED)
    add_panel(slide, 0.75, 4.05, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "特征空间视角", 1.0, 4.32, 1.4, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "理想状态下，同类缺陷应聚在一起；但域偏移让特征按产品线分离，缺陷语义被淹没。", 2.4, 4.28, 9.5, 0.55, font_size = 15, color = INK)
    add_textbox(slide, "域适应要做的事情：把目标域特征簇拉到源域特征簇附近，让缺陷判别边界在目标域上也能生效。", 0.85, 6.0, 11.4, 0.42, font_size = 14, color = NAVY, bold = True)

    # ── Slide 5: Paper Context ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "03 · why this paper", "ADDA 要解决什么", 5, total)
    add_card(slide, "无监督域适应的经典设定", "源域有标签 (Xs, Ys)，目标域只有图像 Xt，没有标签。希望模型在目标域上也能正确分类。", 0.75, 1.75, 3.7, 1.65, accent = BLUE)
    add_card(slide, "对抗学习对齐特征分布", "训练一个域判别器区分源/目标特征，再训练目标编码器欺骗判别器，从而拉近两个域的特征分布。", 4.8, 1.75, 3.7, 1.65, accent = GREEN)
    add_card(slide, "判别式而非生成式", "不生成目标域图像，只对齐特征空间。更轻量、更稳定，尤其适合分类任务。", 8.85, 1.75, 3.7, 1.65, accent = ORANGE)
    add_panel(slide, 0.75, 3.85, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "对应电池场景", 1.0, 4.12, 1.4, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "前五款电池作为源域，第六款无标签图像作为目标域。目标是让第六款的表面纹理被拉到已有特征空间里，源域缺陷分类头可以继续使用。", 2.4, 4.08, 9.5, 0.55, font_size = 15, color = INK)

    # ── Slide 6: Unified Framework ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "04 · framework", "统一对抗域适应框架", 6, total)
    rows = [
        ["设计选择", "DANN (梯度反转)", "Domain Confusion", "CoGAN", "ADDA"],
        ["基础模型", "预训练", "预训练", "预训练", "预训练"],
        ["权重共享", "全部共享", "部分共享", "不共享", "不共享"],
        ["对抗损失", "梯度反转层", "域混淆损失", "GAN loss", "GAN loss"],
        ["建模方式", "判别式", "判别式", "生成式", "判别式"]
    ]
    add_table(slide, rows, 0.75, 1.65, 11.8, 2.8, col_widths = [1.4, 1.8, 1.8, 1.5, 1.5])
    add_panel(slide, 0.75, 4.85, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "关键设计", 1.0, 5.12, 1.2, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "权重不共享 + GAN loss：目标编码器可以自由适应目标域的低层纹理，同时被约束在对齐的特征空间里。", 2.2, 5.08, 9.6, 0.55, font_size = 15, color = INK)

    # ── Slide 7: Three-Stage Pipeline ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "05 · pipeline", "三阶段训练流程", 7, total)
    stages = [
        ("1", "源域预训练", "用源域标签训练源编码器 Ms 和分类器 C。\n目标：最小化分类损失 L_cls。\n此阶段完成后，Ms 和 C 参数固定。", NAVY),
        ("2", "对抗对齐", "固定 Ms，训练域判别器 D 和目标编码器 Mt。\nD 学习区分源/目标特征；Mt 学习欺骗 D。\n对抗博弈让 Mt(Xt) ≈ Ms(Xs)。", BLUE),
        ("3", "目标域推理", "测试时只用目标编码器 + 源分类器。\n输入目标图像 → Mt 提取特征 → C 输出分类。\n不需要判别器 D，也不需要目标域标签。", GREEN)
    ]
    for idx, (num, title, body, accent) in enumerate(stages):
        x = 0.75 + idx * 4.05
        add_panel(slide, x, 1.75, 3.65, 3.8)
        add_badge(slide, num, x + 0.18, 2.0, 0.5, 0.4, bg_color = accent)
        add_textbox(slide, title, x + 0.18, 2.65, 3.2, 0.4, font_size = 18, color = NAVY, bold = True)
        add_textbox(slide, body, x + 0.18, 3.2, 3.2, 2.0, font_size = 13, color = MUTED)
    add_textbox(slide, "对应电池场景：阶段1 = 前五款训练缺陷分类器；阶段2 = 用第六款无标签图像做对抗对齐；阶段3 = 第六款图像直接进分类器出结果。", 0.85, 6.0, 11.4, 0.42, font_size = 13.5, color = NAVY, bold = True)

    # ── Slide 8: Architecture Diagram ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "06 · architecture", "ADDA 架构图", 8, total)
    add_panel(slide, 0.75, 1.55, 11.8, 3.5)
    add_image_fit(slide, base / "adda_pipeline_figure3.png", 0.95, 1.72, 11.4, 3.15)
    add_card(slide, "关键观察", "Source CNN 训练后固定（虚线），Target CNN 单独学习自己的低层映射。域判别器 D 只在第二阶段参与。", 0.75, 5.45, 5.7, 1.2, accent = BLUE)
    add_card(slide, "对应电池场景", "Source CNN = 前五款特征提取器；Target CNN = 第六款特征提取器。分类器 C 复用前五款的缺陷分类头。", 6.8, 5.45, 5.75, 1.2, accent = GREEN)

    # ── Slide 9: Objective Function ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "07 · math", "三个优化目标：分类 + 判别 + 欺骗", 9, total)
    add_panel(slide, 0.75, 1.55, 6.8, 4.5)
    add_image_fit(slide, base / "adda_objective_eq9.png", 0.95, 1.72, 6.4, 4.15)
    formulas = [
        ("① 分类损失", "min L_cls(Xs, Ys)\n用源域标签训练源编码器和分类器", BLUE),
        ("② 判别器损失", "min L_advD\n让 D 正确区分源/目标特征", ORANGE),
        ("③ 目标编码器损失", "min L_advM\n让 Mt 特征被 D 误认为源域特征", GREEN)
    ]
    for idx, (title, body, accent) in enumerate(formulas):
        y = 1.65 + idx * 1.55
        add_card(slide, title, body, 7.85, y, 4.7, 1.2, accent = accent)

    # ── Slide 10: Why It Works ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "08 · intuition", "为什么对抗对齐有效", 10, total)
    add_card(slide, "对齐前", "源域和目标域特征分布差异大\n源域分类边界在目标域上失效\n模型预测塌缩：全部预测为多数类", 0.75, 1.75, 5.7, 1.8, accent = RED)
    add_card(slide, "对齐后", "目标域特征被拉到源域特征空间\n源域分类边界可以继续使用\n预测分布更分散，更接近有监督结果", 7.35, 1.75, 5.7, 1.8, accent = GREEN)
    add_panel(slide, 0.75, 4.05, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "核心逻辑", 1.0, 4.32, 1.2, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "如果 Mt(Xt) ≈ Ms(Xs)，那么源域学到的分类器 C 就能直接用到目标域特征上。", 2.2, 4.28, 9.6, 0.55, font_size = 15, color = INK)
    add_textbox(slide, "电池场景含义：如果第六款的特征被拉到前五款的特征空间里，前五款训练的划痕分类器就可以直接处理第六款图像。", 0.85, 5.85, 11.4, 0.42, font_size = 14, color = NAVY, bold = True)

    # ── Slide 11: Weight Sharing ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "09 · design choice", "为什么不共享权重", 11, total)
    rows = [
        ["维度", "共享权重 (DANN)", "不共享权重 (ADDA)"],
        ["底层映射", "源/目标用同一套，被绑死", "目标编码器独立学习"],
        ["大域偏移", "可能限制对齐效果", "SVHN→MNIST 明显更稳定"],
        ["约束作用", "全层受到梯度反转约束", "只在高层特征空间对齐"],
        ["代表方法", "梯度反转、Domain Confusion", "ADDA (GAN loss)"]
    ]
    add_table(slide, rows, 0.75, 1.65, 11.8, 2.8, col_widths = [1.3, 2.5, 2.5])
    add_panel(slide, 0.75, 4.85, 11.8, 1.35, fill = LIGHT)
    add_textbox(slide, "电池场景意义", 1.0, 5.12, 1.4, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "不同电池型号的表面纹理差异可能很大——光泽度、颗粒感、反光模式都不同。让目标编码器独立学习底层映射，可以更好地捕捉第六款特有的纹理特征。", 2.4, 5.05, 9.5, 0.55, font_size = 14, color = INK)

    # ── Slide 12: Prerequisites ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "10 · prerequisite", "ADDA 的前提条件：不是万能的", 12, total)
    add_card(slide, "✓ 前提 1：目标域有无标签图像", "ADDA 必须有目标域图像 Xt 来训练目标编码器和判别器。如果第六款电池连无标签图像都没有，ADDA 无法启动。", 0.75, 1.75, 3.7, 1.65, accent = BLUE)
    add_card(slide, "? 前提 2：目标域图像包含缺陷样本", "ADDA 对齐的是边缘分布 P(X)，不是条件分布 P(Y|X)。如果目标域只有正常样本，划痕子空间可能无法被正确对齐。", 4.8, 1.75, 3.7, 1.65, accent = ORANGE)
    add_card(slide, "! 前提 3：源/目标类别空间一致", "ADDA 假设源域和目标域共享相同的类别标签集合。如果第六款出现了前五款没有的缺陷类型，ADDA 无法处理。", 8.85, 1.75, 3.7, 1.65, accent = RED)
    add_panel(slide, 0.75, 3.85, 11.8, 1.55, fill = LIGHT)
    add_textbox(slide, "风险提醒", 1.0, 4.12, 1.2, 0.3, font_size = 14, color = RED, bold = True)
    add_textbox(slide, "如果第六款只有正常样本，ADDA 可能只对齐正常表面纹理，无法保证划痕子空间被正确迁移。", 2.2, 4.08, 9.6, 0.55, font_size = 15, color = INK)
    add_bullets(slide, [
        "ADDA 对齐的是边缘分布 P(X)，不是条件分布 P(Y|X) —— 这是核心局限",
        "对电池场景的意义：前五款的划痕标签不能自动覆盖第六款的划痕特征"
    ], 0.85, 5.85, 11.4, 0.8, font_size = 13)

    # ── Slide 13: Digits Results ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "11 · evidence", "实验 1：数字数据集跨域迁移", 13, total)
    add_panel(slide, 0.65, 1.55, 4.05, 4.7)
    metrics = [
        ("90.1%", "USPS → MNIST (ADDA)", BLUE),
        ("89.1%", "USPS → MNIST (CoGAN)", MUTED),
        ("76.0%", "SVHN → MNIST (ADDA)", BLUE),
        ("—", "SVHN → MNIST (CoGAN 未收敛)", RED)
    ]
    for idx, (val, label, accent) in enumerate(metrics):
        add_textbox(slide, val, 0.95, 1.85 + idx * 1.0, 1.3, 0.4, font_size = 22, color = accent, bold = True)
        add_textbox(slide, label, 2.25, 1.92 + idx * 1.0, 2.05, 0.32, font_size = 10.5, color = MUTED)
    add_panel(slide, 5.0, 1.55, 7.55, 4.7)
    add_image_fit(slide, base / "digits_results_table2.png", 5.2, 1.72, 7.15, 4.35)
    add_textbox(slide, "域偏移越大（SVHN → MNIST），判别式方法越稳定；生成式方法 (CoGAN) 在大偏移下可能不收敛。", 0.9, 6.43, 11.4, 0.36, font_size = 13, color = NAVY, bold = True)

    # ── Slide 14: NYUD Results ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "12 · cross-modal", "实验 2：跨模态目标分类", 14, total)
    add_card(slide, "从 RGB 到深度图：域偏移极大", "RGB 图像和深度/HHA 图的成像原理完全不同，是最接近工业跨域的场景。", 0.75, 1.75, 5.7, 1.35, accent = BLUE)
    add_card(slide, "Source only: 13.9% → ADDA: 21.1%", "整体准确率提升显著。无目标标签时，对抗对齐带来了 7 个百分点以上的提升。", 6.85, 1.75, 5.7, 1.35, accent = GREEN)
    add_panel(slide, 0.75, 3.55, 11.8, 3.0)
    add_image_fit(slide, base / "nyud_results_table3.png", 0.95, 3.72, 11.4, 2.65)
    add_textbox(slide, "并非所有类别都提升——某些类别在域适应后仍然无法恢复。这对缺陷检测是重要警示。", 0.85, 6.85, 11.4, 0.36, font_size = 13, color = NAVY, bold = True)

    # ── Slide 15: Confusion Matrix ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "13 · qualitative", "预测塌缩 vs 对齐后恢复", 15, total)
    add_panel(slide, 0.75, 1.55, 11.8, 3.2)
    add_image_fit(slide, base / "confusion_figure5.png", 0.95, 1.72, 11.4, 2.85)
    add_card(slide, "Source only 的症状", "目标域预测严重塌缩。大部分样本被预测为同一个多数类。混淆矩阵呈现“一列高、其余低”。", 0.75, 5.15, 5.7, 1.35, accent = RED)
    add_card(slide, "ADDA 对齐后", "预测分布更分散、更均匀。更接近有目标域监督训练的结果。但仍有部分类别恢复不完全。", 6.85, 5.15, 5.7, 1.35, accent = GREEN)

    # ── Slide 16: Application Analysis ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "14 · application", "能否用到电池缺陷检测？", 16, total)
    add_card(slide, "✓ 可以用作域适应基线", "前五款作源域、第六款正常图像作目标域，ADDA 可以削弱产品外观域偏移。", 0.75, 1.75, 5.7, 1.35, accent = GREEN)
    add_card(slide, "✓ 不需要第六款标签", "只需要第六款的无标签图像就能启动对抗对齐。产线正常品图像即可。", 6.85, 1.75, 5.7, 1.35, accent = BLUE)
    add_card(slide, "! 不能单独解决划痕迁移", "ADDA 对齐边缘分布，不保证划痕子空间被正确对齐。如果第六款无标签数据中没有划痕，对抗对齐看不到目标域划痕的真实形态。", 0.75, 3.5, 11.8, 1.5, accent = RED)
    rows = [
        ["", "ADDA 做到的", "ADDA 做不到的"],
        ["对齐目标", "Ms(Xs) ≈ Mt(Xt)\n边缘特征分布对齐", "Ps(Y|Ms(Xs)) ≈ Pt(Y|Mt(Xt))\n类别条件分布对齐"]
    ]
    add_table(slide, rows, 0.75, 5.4, 11.8, 1.2, col_widths = [1.3, 2.5, 2.5])

    # ── Slide 17: Combined Solution ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "15 · proposed solution", "更稳妥的组合方案", 17, total)
    steps = [
        ("1", "多源缺陷库", "五款电池作为多源域，提供更丰富的缺陷特征覆盖。", NAVY),
        ("2", "ADDA 域适应", "用对抗对齐削弱产品外观域偏移，作为域适应基线。", BLUE),
        ("3", "类条件约束", "加入原型约束或 MMD/CORAL，保证划痕子空间也正确对齐。", GREEN),
        ("4", "少量主动标注", "对第六款少量高不确定性样本做人工标注，弥补条件分布对齐不足。", ORANGE),
        ("5", "异常检测兜底", "用正常样本 PCA 残差或 Mahalanobis 距离做异常检测，覆盖未知的缺陷类型。", RED)
    ]
    for idx, (num, title, body, accent) in enumerate(steps):
        x = 0.6 + idx * 2.45
        add_panel(slide, x, 1.75, 2.2, 3.8)
        add_badge(slide, num, x + 0.15, 2.0, 0.45, 0.36, bg_color = accent)
        add_textbox(slide, title, x + 0.15, 2.55, 1.85, 0.4, font_size = 15, color = NAVY, bold = True)
        add_textbox(slide, body, x + 0.15, 3.1, 1.85, 2.0, font_size = 11.5, color = MUTED)
    add_panel(slide, 0.75, 6.0, 11.8, 0.75, fill = LIGHT)
    add_textbox(slide, "核心策略", 1.0, 6.18, 1.2, 0.3, font_size = 13, color = NAVY, bold = True)
    add_textbox(slide, "ADDA 做域适应基础层，类条件约束保缺陷子空间，少量标注补条件分布，异常检测兜底未知类。", 2.2, 6.18, 9.6, 0.35, font_size = 14, color = INK)

    # ── Slide 18: Summary ──
    slide = prs.slides.add_slide(blank)
    add_header(slide, "16 · summary", "ADDA 的定位与边界", 18, total)
    add_card(slide, "ADDA 做对了什么", "清晰的无监督域适应框架：源域预训练 → 对抗对齐 → 目标推理\n判别式方法比生成式更稳定，尤其是大域偏移下\n不共享权重设计允许目标域有独立低层映射\n作为域适应基线，方法简洁、实验充分", 0.75, 1.75, 5.7, 2.4, accent = GREEN)
    add_card(slide, "ADDA 的局限", "只对齐边缘分布，不保证类别条件对齐\n目标域无标签数据中必须有缺陷样本覆盖\n不能处理开放集或新类别场景\n不适合单独作为工业缺陷迁移的最终方案", 6.85, 1.75, 5.7, 2.4, accent = RED)
    add_panel(slide, 0.75, 4.55, 11.8, 1.55, fill = LIGHT)
    add_textbox(slide, "结论", 1.0, 4.82, 0.8, 0.3, font_size = 14, color = NAVY, bold = True)
    add_textbox(slide, "ADDA 适合做域适应基线——削弱产品外观偏移；但划痕子空间对齐需要额外约束（类条件 + 少量标注 + 异常检测）。", 1.8, 4.78, 10.1, 0.55, font_size = 15, color = INK)
    add_bullets(slide, [
        "如果有问题，我们可以讨论 ADDA 和 DANN 的选择、类条件约束的具体实现、或主动标注的策略设计。"
    ], 0.85, 6.5, 11.4, 0.5, font_size = 13)

    prs.save(output_path)


if __name__ == "__main__":
    build_deck("2026-04-28/results/html_ppt/ADDA_paper_report.pptx")
